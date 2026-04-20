"""
Genereer een PowerPoint presentatie (audit_presentatie.pptx) op basis van
de inhoud van de management samenvatting audit 2026.

Gebruik:
  python3 -m audit.make_pptx
"""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Kleurenpalet (zelfde als de HTML)
# ---------------------------------------------------------------------------
DONKER_BLAUW   = RGBColor(0x1a, 0x1a, 0x2e)
MIDDEL_BLAUW   = RGBColor(0x16, 0x21, 0x3e)
ACCENT_BLAUW   = RGBColor(0x0f, 0x34, 0x60)
ROOD           = RGBColor(0xe9, 0x45, 0x60)
ORANJE         = RGBColor(0xf0, 0xa5, 0x00)
GROEN          = RGBColor(0x2e, 0xcc, 0x71)
GRIJS          = RGBColor(0x99, 0x99, 0x99)
LICHT_GRIJS    = RGBColor(0xf8, 0xf9, 0xfa)
TEKST          = RGBColor(0x33, 0x33, 0x33)
WIT            = RGBColor(0xff, 0xff, 0xff)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _nieuwe_slide(prs: Presentation, layout_idx: int = 6):
    """Lege slide (blank layout)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _achtergrond(slide, kleur: RGBColor = WIT):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = kleur


def _textbox(slide, left, top, width, height,
             tekst: str, grootte: int, vet: bool = False,
             kleur: RGBColor = TEKST, align=PP_ALIGN.LEFT,
             wrap: bool = True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = tekst
    run.font.size = Pt(grootte)
    run.font.bold = vet
    run.font.color.rgb = kleur


def _rode_lijn(slide, top):
    """Horizontale rode scheidingslijn onder de titel."""
    from pptx.util import Pt as _Pt
    lijn = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), top, Inches(12.33), Inches(0.05),
    )
    lijn.fill.solid()
    lijn.fill.fore_color.rgb = ROOD
    lijn.line.fill.background()


def _kader(slide, left, top, width, height,
           randkleur: RGBColor = ROOD,
           achtergrond: RGBColor = LICHT_GRIJS):
    """Gekleurde border-left kader (gesimuleerd als twee rechthoeken)."""
    # Linkse balk
    balk = slide.shapes.add_shape(1, left, top, Inches(0.07), height)
    balk.fill.solid()
    balk.fill.fore_color.rgb = randkleur
    balk.line.fill.background()
    # Achtergrondvlak
    vlak = slide.shapes.add_shape(
        1, left + Inches(0.07), top, width - Inches(0.07), height,
    )
    vlak.fill.solid()
    vlak.fill.fore_color.rgb = achtergrond
    vlak.line.fill.background()


def _kader_tekst(slide, left, top, width, height,
                 inhoud: list[tuple[str, bool, int]],
                 randkleur: RGBColor = ROOD):
    """Kader + meerdere tekstregels [(tekst, vet, grootte)]."""
    _kader(slide, left, top, width, height, randkleur=randkleur)
    txb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.08),
        width - Inches(0.25), height - Inches(0.16),
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (tekst, vet, grootte) in enumerate(inhoud):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = tekst
        run.font.size = Pt(grootte)
        run.font.bold = vet
        run.font.color.rgb = TEKST


def _bullet_lijst(slide, left, top, width, height,
                  items: list[str], grootte: int = 16):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"→  {item}"
        run.font.size = Pt(grootte)
        run.font.color.rgb = TEKST


def _stat_blok(slide, left, top, getal: str, label: str, sublabel: str,
               kleur: RGBColor):
    _textbox(slide, left, top, Inches(2.2), Inches(0.9),
             getal, 40, vet=True, kleur=kleur, align=PP_ALIGN.CENTER)
    _textbox(slide, left, top + Inches(0.85), Inches(2.2), Inches(0.4),
             label, 14, vet=True, kleur=kleur, align=PP_ALIGN.CENTER)
    _textbox(slide, left, top + Inches(1.2), Inches(2.2), Inches(0.35),
             sublabel, 11, kleur=GRIJS, align=PP_ALIGN.CENTER)


def _slide_titel(slide, tekst: str):
    _textbox(slide, Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.65),
             tekst, 28, vet=True, kleur=DONKER_BLAUW)
    _rode_lijn(slide, Inches(0.9))


def _footer(slide):
    _textbox(slide,
             Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.3),
             "ISO Audit 2026  ·  Conduction  ·  Vertrouwelijk",
             9, kleur=GRIJS, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_titel(prs):
    s = _nieuwe_slide(prs)
    _achtergrond(s, DONKER_BLAUW)

    _textbox(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.3),
             "ISO 9001 + 27001", 48, vet=True, kleur=WIT, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "Interne Audit 2026", 36, vet=False, kleur=ROOD, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(3.9), Inches(10), Inches(0.6),
             "Management Samenvatting", 22, kleur=RGBColor(0xcc, 0xcc, 0xcc),
             align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             "Conduction  ·  24 maart 2026  ·  Vertrouwelijk",
             13, kleur=GRIJS, align=PP_ALIGN.CENTER)


def slide_organisatie(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Organisatie in Beweging")

    items = [
        "Conduction is de afgelopen jaren significant gegroeid en veranderd",
        "Documentatie en procedures zijn niet altijd meegegroeid",
        "208 van de 282 documenten zijn ouder dan 2 jaar — zij weerspiegelen niet altijd meer de huidige situatie",
        "Deze audit meet de huidige staat: alleen documenten uit 2024+ zijn beoordeeld",
    ]
    _bullet_lijst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(2.5), items, 17)

    _kader_tekst(s, Inches(0.5), Inches(3.8), Inches(12.33), Inches(1.4),
                 [("Aanbeveling: Stel een jaarlijkse documentatierevisieprocedure in. "
                   "Verouderde documenten zijn geen bewijs — ze zijn een risico "
                   "(bijv. documenten Mat + geen vertrouwelijkheidsniveau).", False, 15)],
                 randkleur=ROOD)
    _footer(s)


def slide_aanpak(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Aanpak: AI-ondersteunde Audit")

    items = [
        "Deze audit is mede-uitgevoerd met geautomatiseerde AI-analyse (Claude)",
        "Hiervoor is gebruik gemaakt op de doelstellingen mede te behalen (automatiseren)",
        "Alle Drive-documenten én Miro-borden zijn systematisch gescand",
        "Doel: de blinde vlek van de auditor die meerdere petjes draagt elimineren",
        "Bevindingen zijn gegenereerd op basis van daadwerkelijk aanwezige bewijslast — niet op aannames",
    ]
    _bullet_lijst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(2.8), items, 17)

    _kader_tekst(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(1.1),
                 [("De organisatie hanteert \"afwijking\" als term voor non-conformiteit "
                   "en heeft hiervoor gedocumenteerde procedures. Dit is meegewogen in de beoordeling.", False, 15)],
                 randkleur=ACCENT_BLAUW)
    _footer(s)


def slide_cijfers(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Resultaten in Één Oogopslag")

    _stat_blok(s, Inches(0.8),  Inches(1.4), "68",  "NC",           "non-conformiteiten", ROOD)
    _stat_blok(s, Inches(3.3),  Inches(1.4), "288", "OFI",          "verbeterpunten",      ORANJE)
    _stat_blok(s, Inches(5.8),  Inches(1.4), "93",  "✓ positief",   "positief",            GROEN)
    _stat_blok(s, Inches(8.3),  Inches(1.4), "208", "gearchiveerd", "te herzien (>2 jaar)", GRIJS)

    _textbox(s, Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.4),
             "Beoordeeld: 68 actieve documenten + 170 Miro-notities",
             13, kleur=GRIJS, align=PP_ALIGN.CENTER)
    _footer(s)


def slide_kritiek(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Kritieke Bevindingen  [ NC ]")

    _kader_tekst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(1.5),
                 [("5.11 — Retournering van activa", True, 16),
                  ("Geen gedocumenteerde procedure voor teruggave van bedrijfsmiddelen bij vertrek of klantbeëindiging. "
                   "Actie: offboarding-checklist formaliseren (informele procedure bestaat al).", False, 14)],
                 randkleur=ROOD)

    _kader_tekst(s, Inches(0.5), Inches(2.8), Inches(12.33), Inches(1.3),
                 [("5.12 / 5.13 — Informatieclassificatie & -labeling", True, 16),
                  ("Geen classificatieschema (vertrouwelijk / intern / openbaar) of labelingprocedure gedocumenteerd.", False, 14)],
                 randkleur=ROOD)

    _textbox(s, Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.5),
             "De meeste NC's betreffen ontbrekende documentatie van bestaande praktijken — niet het ontbreken van de praktijk zelf.",
             13, kleur=GRIJS)
    _footer(s)


def slide_plan_nc(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Plan van Aanpak NC's")

    items = [
        "Q2 2026: Informatieclassificatiebeleid opstellen (5.12 / 5.13)",
        "Q3 2026: Documentatierevisieprocedure invoeren — jaarlijkse review cyclus "
        "niet enkel actualiseren, maar ook opruimen en check op oud-medewerkers",
    ]
    _bullet_lijst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(1.6), items, 17)

    _kader_tekst(s, Inches(0.5), Inches(3.0), Inches(12.33), Inches(1.8),
                 [("68 NC's zijn minor en mede opgesteld door AI. De meeste zijn documentatiegaps, "
                   "geen procesfouten. Met een gerichte sprint zijn ze binnen één kwartaal gesloten. "
                   "Mark doet extra controle en maakt issues aan waar nodig en zet deze op naam.", False, 15)],
                 randkleur=ROOD)
    _footer(s)


def slide_ofi(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Kansen voor Verbetering  [ OFI ]")

    _kader_tekst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(1.3),
                 [("9.1 — Monitoring & meting  (12×)", True, 15),
                  ("Q3/Q4 interne audit onvolledig door personeelswisseling. "
                   "Structurele auditagenda met backup-auditor gewenst.", False, 13)],
                 randkleur=ORANJE)

    _kader_tekst(s, Inches(0.5), Inches(2.55), Inches(12.33), Inches(1.5),
                 [("5.31 — Wet- en regelgeving  (18×)", True, 15),
                  ("Actiepunten voor compliance-update aanwezig maar niet opgevolgd. "
                   "Eigenaar en deadline toewijzen. Aanbeveling: zet om naar NC, OFI en positief.", False, 13)],
                 randkleur=ORANJE)

    _kader_tekst(s, Inches(0.5), Inches(4.2), Inches(12.33), Inches(1.4),
                 [("8.16 — Monitoring  (21×)", True, 15),
                  ("Camerabewaking en verwerkingsregister benoemd maar niet volledig ingericht. "
                   "Monitoring gebeurt nu nog veel door de mens — scope niet expliciet gedocumenteerd.", False, 13)],
                 randkleur=ORANJE)
    _footer(s)


def slide_positief(prs):
    s = _nieuwe_slide(prs)
    _slide_titel(s, "Sterke Punten  [ ✓ positief ]")

    items = [
        "10.2 — Afwijkingsprocedure werkt: gedocumenteerde memo's tonen actief NC-beheer",
        "Interne audits worden structureel uitgevoerd (Q1/Q2 2025 compleet)",
        "Incident management: incidenten worden gelogd en opgevolgd",
        "Gecertificeerd: de organisatie voldoet aan de kern van beide normen",
    ]
    _bullet_lijst(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(2.2), items, 17)

    _kader_tekst(s, Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.4),
                 [("De certificering is terecht. De audit toont een organisatie die wil verbeteren "
                   "en de juiste processen in gang heeft gezet. De volgende stap is documentatie bijbenen.", False, 16)],
                 randkleur=GROEN)
    _footer(s)


def slide_conclusie(prs):
    s = _nieuwe_slide(prs)
    _achtergrond(s, LICHT_GRIJS)

    _textbox(s, Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.9),
             "Conclusie", 34, vet=True, kleur=DONKER_BLAUW, align=PP_ALIGN.CENTER)
    _rode_lijn(s, Inches(1.25))

    _textbox(s, Inches(0.8), Inches(1.5), Inches(11.73), Inches(1.0),
             "Conduction is een organisatie in beweging. De audit laat zien dat de processen werken, "
             "maar de documentatie achterblijft. Dat is een kans, geen crisis.",
             18, kleur=TEKST, align=PP_ALIGN.CENTER)

    items = [
        "68 NC's → behapbaar actieplan Q2/Q3 2026",
        "288 OFI's → continue verbeteragenda",
        "93 positieve bevindingen → stevige basis",
    ]
    _bullet_lijst(s, Inches(3.0), Inches(2.8), Inches(7.33), Inches(1.6), items, 18)

    _textbox(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.4),
             "Volledig rapport: Auditrapport_beide_2026-03-24_s05.md",
             12, kleur=GRIJS, align=PP_ALIGN.CENTER)
    _footer(s)


# ---------------------------------------------------------------------------
# Hoofdfunctie
# ---------------------------------------------------------------------------

def maak_presentatie(uitvoer: Path | None = None) -> Path:
    if uitvoer is None:
        uitvoer = Path("output/audit_presentatie_2026-03-24.pptx")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_titel(prs)
    slide_organisatie(prs)
    slide_aanpak(prs)
    slide_cijfers(prs)
    slide_kritiek(prs)
    slide_plan_nc(prs)
    slide_ofi(prs)
    slide_positief(prs)
    slide_conclusie(prs)

    uitvoer.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(uitvoer))
    return uitvoer


if __name__ == "__main__":
    import logging
    logging.basicConfig(level=logging.INFO, format="%(message)s")
    pad = maak_presentatie()
    print(f"Presentatie opgeslagen: {pad}")
